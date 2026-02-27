# scripts/generate_pptx_from_baseline.py

from __future__ import annotations

import re
from typing import Dict, Any, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from ml.models.baseline import generate_deck
from ml.prompt_templates import PitchDeckPromptConfig


# Краткое описание стартапа для генерации примера
BRIEF = (
    "AI Pitch Deck Generator is a SaaS tool that helps startup founders "
    "generate investor-ready pitch decks from short textual briefs."
)


# ---------- график для Market & Financials ----------

def extract_numbers(text: str) -> List[float]:
    """Пробуем вытащить числа из текста (для графика)."""
    nums = re.findall(r"\d+\.?\d*", text)
    return [float(n) for n in nums]


def add_market_chart(slide, bullets: List[str]) -> None:
    """
    Столбчатый график для слайда Market & Financials.
    Если нашли числа в тексте — используем их,
    иначе рисуем условный TAM/SAM/SOM.
    """
    joined = " ".join(bullets)
    numbers = extract_numbers(joined)

    if len(numbers) >= 3:
        categories = [f"Metric {i+1}" for i in range(len(numbers))]
        values = numbers
    else:
        categories = ["TAM", "SAM", "SOM"]
        values = (100, 60, 30)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Market/Financials", values)

    x = Inches(7.0)
    y = Inches(2.0)
    cx = Inches(3.3)
    cy = Inches(3.0)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = chart_frame.chart

    chart.has_legend = False
    chart.value_axis.has_title = False
    chart.category_axis.has_title = False


# ---------- генерация презентации ----------

def build_presentation(deck: Dict[str, Any], filename: str = "pitch_deck_generated.pptx") -> None:
    prs = Presentation()

    # Титульный слайд
    title_layout = prs.slide_layouts[0]
    slide0 = prs.slides.add_slide(title_layout)
    slide0.shapes.title.text = "AI Pitch Deck Generator"
    slide0.placeholders[1].text = "Auto-generated investor deck (baseline model)"

    # Контентные слайды
    bullet_layout = prs.slide_layouts[1]

    for s in deck["slides"]:
        section = s.get("section", "")
        title = s.get("title", "")
        bullets = s.get("bullets") or []

        slide = prs.slides.add_slide(bullet_layout)

        # Заголовок
        title_shape = slide.shapes.title
        title_shape.text = f"{section}: {title}"
        title_tf = title_shape.text_frame
        if title_tf.paragraphs:
            title_tf.paragraphs[0].font.size = Pt(30)

        # Текстовый блок
        body = slide.placeholders[1]
        # Немного смещаем и ограничиваем ширину,
        # чтобы на Market & Financials справа было место под график
        body.left = Inches(0.7)
        body.top = Inches(2.0)
        body.width = Inches(5.8)
        body.height = Inches(4.0)

        tf = body.text_frame
        tf.clear()

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

        # Только для Market & Financials добавляем график справа,
        # для остальных слайдов — чистый текст.
        if "market" in section.lower() and "financial" in section.lower():
            add_market_chart(slide, bullets)

    prs.save(filename)
    print(f"Presentation saved to: {filename}")


def main() -> None:
    config = PitchDeckPromptConfig(
        target_slide_count=10,
        tone="formal",
        audience="early-stage VC",
        language="en",
    )
    deck = generate_deck(BRIEF, config)
    build_presentation(deck, "pitch_deck_generated.pptx")


if __name__ == "__main__":
    main()
